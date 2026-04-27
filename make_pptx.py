"""PPTX iperprofessionali pitch-grade. No costi, no brand specifici."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

DOCS = r"C:\Users\Giovanni\graficanappa-gestionale\docs"

# Brand palette consulting-grade
NAVY = RGBColor(0x0A, 0x1F, 0x44)
NAVY_DEEP = RGBColor(0x06, 0x14, 0x2E)
BLUE = RGBColor(0x2B, 0x6C, 0xFF)
TEAL = RGBColor(0x0E, 0x9F, 0x9F)
INK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD4, 0xA8, 0x4B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)


def new_pres():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w:
            s.line.width = line_w
    if hasattr(s, "text_frame"):
        s.text_frame.margin_left = Inches(0.1)
        s.text_frame.margin_right = Inches(0.1)
        s.text_frame.margin_top = Inches(0.05)
        s.text_frame.margin_bottom = Inches(0.05)
    return s


def text(slide, x, y, w, h, t, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri", spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = t
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=INK, gap=8, marker_color=None, marker_char="—"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        m = p.add_run()
        m.text = marker_char + "  "
        m.font.size = Pt(size)
        m.font.bold = True
        m.font.color.rgb = marker_color or BLUE
        m.font.name = "Calibri"
        r = p.add_run()
        r.text = it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def header(slide, eyebrow, title, sub=None, num=None, monogram=None):
    rect(slide, 0, 0, Inches(13.333), Inches(0.06), BLUE)
    # monogram top-left
    if monogram:
        rect(slide, Inches(0.6), Inches(0.25), Inches(0.5), Inches(0.5), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(slide, Inches(0.6), Inches(0.25), Inches(0.5), Inches(0.5), monogram, size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        eyebrow_x = Inches(1.25)
    else:
        eyebrow_x = Inches(0.6)
    if num:
        text(slide, eyebrow_x, Inches(0.32), Inches(0.7), Inches(0.35), num, size=10, bold=True, color=MUTED)
        rect(slide, eyebrow_x + Inches(0.7), Inches(0.46), Inches(0.04), Inches(0.04), BLUE, shape=MSO_SHAPE.OVAL)
        text(slide, eyebrow_x + Inches(0.85), Inches(0.32), Inches(11), Inches(0.35), eyebrow.upper(), size=10, bold=True, color=BLUE)
    else:
        text(slide, eyebrow_x, Inches(0.32), Inches(12), Inches(0.35), eyebrow.upper(), size=10, bold=True, color=BLUE)
    text(slide, Inches(0.6), Inches(0.78), Inches(12), Inches(0.7), title, size=26, bold=True, color=INK)
    if sub:
        text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), sub, size=13, color=SLATE)
    rect(slide, Inches(0.6), Inches(1.75), Inches(12.13), Inches(0.015), LINE)


def footer(slide, num, total, brand="Confidenziale"):
    text(slide, Inches(0.6), Inches(7.1), Inches(8), Inches(0.3), brand, size=10, color=MUTED)
    # page num refined: "01" · BLUE dot · "17"
    text(slide, Inches(11.5), Inches(7.1), Inches(0.5), Inches(0.3), f"{num:02d}", size=10, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    rect(slide, Inches(12.05), Inches(7.21), Inches(0.06), Inches(0.06), BLUE, shape=MSO_SHAPE.OVAL)
    text(slide, Inches(12.18), Inches(7.1), Inches(0.5), Inches(0.3), f"{total:02d}", size=10, color=MUTED, align=PP_ALIGN.LEFT)


def kpi(slide, x, y, w, h, value, label, color=BLUE, size=36):
    rect(slide, x, y, w, h, WHITE, line=LINE, line_w=Pt(0.75))
    rect(slide, x, y, w, Inches(0.06), color)
    text(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.4), Inches(1.0), value, size=size, bold=True, color=color)
    text(slide, x + Inches(0.3), y + h - Inches(0.6), w - Inches(0.4), Inches(0.5), label, size=11, color=SLATE)


def hero_title(p, eyebrow, title, subtitle, meta):
    s = blank(p)
    rect(s, 0, 0, Inches(13.333), Inches(7.5), NAVY_DEEP)
    rect(s, Inches(10.7), 0, Inches(2.63), Inches(7.5), NAVY)
    rect(s, Inches(12.5), Inches(1.0), Inches(0.04), Inches(5.5), GOLD)
    rect(s, Inches(10.85), Inches(1.0), Inches(0.04), Inches(5.5), BLUE)
    rect(s, 0, 0, Inches(0.5), Inches(0.06), GOLD)
    rect(s, 0, 0, Inches(0.06), Inches(0.5), GOLD)
    rect(s, Inches(0.7), Inches(1.4), Inches(0.6), Inches(0.04), GOLD)
    text(s, Inches(0.7), Inches(1.5), Inches(9.5), Inches(0.4), eyebrow, size=12, bold=True, color=GOLD)
    text(s, Inches(0.7), Inches(2.0), Inches(9.5), Inches(2.8), title, size=46, bold=True, color=WHITE, spacing=1.1)
    text(s, Inches(0.7), Inches(5.2), Inches(9.5), Inches(1.4), subtitle, size=17, color=MUTED, spacing=1.3)
    rect(s, Inches(0.7), Inches(6.85), Inches(0.4), Inches(0.04), GOLD)
    text(s, Inches(0.7), Inches(6.95), Inches(10), Inches(0.4), meta, size=11, color=MUTED)
    return s


def section_divider(p, num, eyebrow, title):
    s = blank(p)
    rect(s, 0, 0, Inches(13.333), Inches(7.5), BG)
    rect(s, 0, 0, Inches(0.5), Inches(7.5), BLUE)
    # dot pattern decorativo top-right
    for row in range(8):
        for col in range(8):
            x = Inches(8.5) + col * Inches(0.18)
            y = Inches(0.6) + row * Inches(0.18)
            rect(s, x, y, Inches(0.04), Inches(0.04), LINE, shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.2), Inches(2.0), Inches(11), Inches(1.5), num, size=120, bold=True, color=LINE, spacing=1.0)
    rect(s, Inches(1.2), Inches(3.7), Inches(0.6), Inches(0.04), BLUE)
    text(s, Inches(1.2), Inches(3.8), Inches(11), Inches(0.4), eyebrow.upper(), size=12, bold=True, color=BLUE)
    text(s, Inches(1.2), Inches(4.2), Inches(11), Inches(2.5), title, size=44, bold=True, color=INK, spacing=1.1)
    # bottom anchor
    rect(s, Inches(1.2), Inches(6.8), Inches(0.4), Inches(0.04), GOLD)
    text(s, Inches(1.2), Inches(6.9), Inches(11), Inches(0.3), "GRAFICA NAPPA", size=9, bold=True, color=MUTED)
    return s


def closing_slide(p, title, sub, contact):
    s = blank(p)
    rect(s, 0, 0, Inches(13.333), Inches(7.5), NAVY_DEEP)
    rect(s, Inches(8.5), 0, Inches(4.83), Inches(7.5), NAVY)
    rect(s, Inches(0.7), Inches(1.6), Inches(0.6), Inches(0.04), GOLD)
    text(s, Inches(0.7), Inches(1.7), Inches(11), Inches(0.4), "PROSSIMI PASSI", size=12, bold=True, color=GOLD)
    text(s, Inches(0.7), Inches(2.2), Inches(11), Inches(2.5), title, size=48, bold=True, color=WHITE, spacing=1.05)
    text(s, Inches(0.7), Inches(5.0), Inches(11), Inches(1.2), sub, size=18, color=MUTED, spacing=1.3)
    rect(s, Inches(0.7), Inches(6.6), Inches(0.4), Inches(0.04), GOLD)
    text(s, Inches(0.7), Inches(6.7), Inches(11), Inches(0.4), contact, size=14, color=GOLD)
    text(s, Inches(0.7), Inches(7.15), Inches(11), Inches(0.3), "Confidenziale", size=10, color=MUTED)
    return s


def node(slide, x, y, w, h, title, sub=None, fill=WHITE, line=LINE, accent=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = rect(slide, x, y, w, h, fill, line=line, line_w=Pt(0.75), shape=shape)
    if accent:
        rect(slide, x, y, Inches(0.08), h, accent)
    text(slide, x + Inches(0.18), y + Inches(0.13), w - Inches(0.28), Inches(0.4), title, size=12, bold=True, color=INK)
    if sub:
        text(slide, x + Inches(0.18), y + Inches(0.5), w - Inches(0.28), h - Inches(0.6), sub, size=10, color=SLATE)
    return sh


def arrow(slide, x1, y1, x2, y2, color=MUTED, weight=1.25):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    ln = c.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return c


def notes(slide, text_content):
    """Aggiunge speaker notes alla slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text_content


def donut_label(slide, cx, cy, r, label_text, color):
    rect(slide, cx - r, cy - r, 2 * r, 2 * r, color, shape=MSO_SHAPE.OVAL)
    inner_r = int(r * 0.66)
    rect(slide, cx - inner_r, cy - inner_r, 2 * inner_r, 2 * inner_r, WHITE, shape=MSO_SHAPE.OVAL)
    text(slide, cx - r, cy - Inches(0.32), 2 * r, Inches(0.7), label_text, size=22, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# PPT 1 — Mossa 37 Scheduler
# ============================================================
def build_scheduler():
    p = new_pres()
    total = 19

    # 1 Hero
    hero_title(p,
        "MOSSA 37  ·  PRODUCTION SCHEDULER",
        "La pianificazione,\nfinalmente, in un secondo.",
        "Scheduler real-time per stampa offset e digitale.\nRicostruisce la sequenza ottimale a ogni evento di reparto.",
        "Documento riservato  ·  Aprile 2026")

    # 2 Indice
    s = blank(p)
    header(s, "Agenda", "Dodici minuti. Cinque sezioni. Una decisione.", num="00")
    items = [
        ("01", "Il costo invisibile dell'Excel", "perché la pianificazione manuale erode marginalità"),
        ("02", "Mossa 37, in profondità", "logica di scheduling, criteri, motore decisionale"),
        ("03", "Impatto operativo", "throughput, setup time, lead time interno"),
        ("04", "Roadmap di adozione", "discovery, pilota, go-live in 4 settimane"),
        ("05", "Demo live", "lo scheduler in funzione, sui vostri dati di esempio"),
    ]
    for i, (n, t, d) in enumerate(items):
        y = Inches(2.0) + i * Inches(0.95)
        rect(s, Inches(0.6), y, Inches(12.13), Inches(0.85), WHITE, line=LINE, line_w=Pt(0.5))
        text(s, Inches(0.9), y + Inches(0.18), Inches(0.8), Inches(0.5), n, size=22, bold=True, color=BLUE)
        text(s, Inches(2.0), y + Inches(0.18), Inches(4.5), Inches(0.5), t, size=18, bold=True, color=INK)
        text(s, Inches(6.7), y + Inches(0.22), Inches(6), Inches(0.5), d, size=13, color=SLATE)
    footer(s, 2, total)

    # 3 Sect 01
    section_divider(p, "01", "Il costo invisibile", "L'Excel non ha un prezzo.\nMa ha un conto.")

    # 4 Problema in numeri (no €)
    s = blank(p)
    header(s, "Anatomia di una giornata-tipo", "Tre numeri che misuriamo in ogni reparto stampa", num="01")
    cy = Inches(3.4)
    r1 = Inches(0.95)
    centers = [
        (Inches(2.5), "2,4h", RED, "scheduling manuale", "tempo medio capo reparto / die"),
        (Inches(6.65), "8 ×", AMBER, "ricalcolo / die", "ogni urgenza propaga su più fasi"),
        (Inches(10.8), "30%", BLUE, "OEE non sfruttato", "attese e setup ridondanti evitabili"),
    ]
    for cx, lbl, col, top, sub in centers:
        donut_label(s, cx, cy, r1, lbl, col)
        text(s, cx - r1, cy + r1 + Inches(0.2), 2 * r1, Inches(0.5), top, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s, cx - r1 - Inches(0.5), cy + r1 + Inches(0.7), 2 * r1 + Inches(1), Inches(0.5), sub, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    rect(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(1.2), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.4), "EFFETTO CASCATA", size=11, bold=True, color=GOLD)
    text(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6), "Una variazione genera setup duplicati, fasi rilavorate, lead time fuori standard, decisioni a memoria.", size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 4, total)
    notes(s, "I numeri sono medie su tipografie 10-50 dipendenti. Chiedi al cliente: 'quanto tempo passa il vostro capo reparto su Excel ogni mattina?'. Lascia rispondere prima di proseguire. Se dicono <2h, complimentati e di' che molti stanno peggio.")

    # 5 Pull quote — editoriale
    s = blank(p)
    header(s, "Voce dal reparto", "Capo reparto stampa, tipografia cartotecnica", num="01", monogram="M37")
    # giant decorative quote mark
    text(s, Inches(0.7), Inches(1.7), Inches(2), Inches(2), "“", size=180, bold=True, color=GOLD, font="Calibri")
    rect(s, Inches(2.4), Inches(2.8), Inches(0.15), Inches(3.0), GOLD)
    text(s, Inches(2.85), Inches(2.6), Inches(10), Inches(2.8),
         "La prima ora la passavo a capire cosa fosse saltato\nnella notte. Adesso entro, leggo il Gantt,\ne so dove intervenire.",
         size=26, color=INK, spacing=1.4)
    rect(s, Inches(2.85), Inches(5.5), Inches(0.4), Inches(0.04), GOLD)
    text(s, Inches(2.85), Inches(5.6), Inches(11), Inches(0.4), "CAPO REPARTO STAMPA · TIPOGRAFIA 28 ADDETTI", size=11, bold=True, color=INK)
    text(s, Inches(2.85), Inches(5.95), Inches(11), Inches(0.4), "Testimonianza raccolta su un sito di produzione attivo", size=11, color=SLATE)
    footer(s, 5, total)
    notes(s, "Pull quote chiave. Pausa lunga dopo la lettura. Lascia che il numero (2 ore) faccia il lavoro. Connetti: 'questa frase ce l'ha detta un capo reparto come voi'.")

    # 6 Sect 02
    section_divider(p, "02", "La risposta", "Decisioni di scheduling\ndelegate a un motore.")

    # 7 Value prop + diff (no brand)
    s = blank(p)
    header(s, "Mossa 37, in profondità", "Una proposizione di valore, tre vantaggi competitivi", num="02")
    rect(s, Inches(0.6), Inches(2.4), Inches(12.13), Inches(2.2), BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2),
         "Acquisisce dati da ERP, macchine offset, digitali e presenze.\nRicompone la sequenza ottimale a ogni evento di reparto.\nDistribuisce un Gantt unico, vivo, a tutta la filiera produttiva.",
         size=20, bold=True, color=INK, spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4), "VANTAGGI COMPETITIVI", size=11, bold=True, color=BLUE)
    diffs = [
        ("Domain-native", "Regole calibrate su flussi reali di stampa cartotecnica, non astratte."),
        ("Decisione trasparente", "Cinque criteri ponderati e ispezionabili. Nessun motore opaco."),
        ("Architettura non invasiva", "Si integra sopra il vostro stack. Nessun rip-and-replace."),
    ]
    cw = Inches(3.95)
    for i, (n, d) in enumerate(diffs):
        x = Inches(0.6) + i * (cw + Inches(0.13))
        rect(s, x, Inches(5.4), cw, Inches(1.6), WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(5.4), Inches(0.08), Inches(1.6), GOLD)
        text(s, x + Inches(0.25), Inches(5.55), cw - Inches(0.4), Inches(0.5), n, size=18, bold=True, color=INK)
        text(s, x + Inches(0.25), Inches(6.05), cw - Inches(0.4), Inches(1.0), d, size=12, color=SLATE)
    footer(s, 7, total)

    # 8 Come funziona (generic)
    s = blank(p)
    header(s, "Architettura funzionale", "Quattro sorgenti · un motore decisionale · tre output · sub-secondo", num="02")

    text(s, Inches(0.7), Inches(1.95), Inches(3), Inches(0.35), "INPUT REAL-TIME", size=10, bold=True, color=BLUE)
    inputs = [("ERP gestionale", "commesse + urgenze"), ("Macchine offset", "tempi reali"), ("Macchine digitali", "coda di stampa"), ("Calendari reparto", "disponibilità · opzionale")]
    for i, (n, d) in enumerate(inputs):
        node(s, Inches(0.7), Inches(2.3) + i * Inches(0.95), Inches(2.7), Inches(0.8), n, sub=d, accent=BLUE)

    rect(s, Inches(4.3), Inches(3.3), Inches(2.5), Inches(2.0), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(4.4), Inches(3.45), Inches(2.3), Inches(0.4), "MOTORE", size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text(s, Inches(4.4), Inches(3.8), Inches(2.3), Inches(0.6), "Mossa 37", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(4.85), Inches(4.4), Inches(1.4), Inches(0.04), GOLD)
    text(s, Inches(4.4), Inches(4.5), Inches(2.3), Inches(0.5), "< 1 secondo", size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text(s, Inches(4.4), Inches(4.95), Inches(2.3), Inches(0.4), "calcolo priorità", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    text(s, Inches(7.8), Inches(1.95), Inches(3), Inches(0.35), "OUTPUT", size=10, bold=True, color=GREEN)
    outs = [("Gantt real-time", "per macchina"), ("Batch ottimizzati", "setup minimi"), ("Alert capo reparto", "colli di bottiglia")]
    for i, (n, d) in enumerate(outs):
        node(s, Inches(7.8), Inches(2.5) + i * Inches(1.0), Inches(2.7), Inches(0.85), n, sub=d, accent=GREEN)

    rect(s, Inches(10.85), Inches(2.4), Inches(1.95), Inches(2.85), BG, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(10.95), Inches(2.6), Inches(1.75), Inches(0.4), "RISULTATO", size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    text(s, Inches(10.95), Inches(3.0), Inches(1.75), Inches(0.7), "50+", size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    text(s, Inches(10.95), Inches(3.7), Inches(1.75), Inches(0.4), "commesse", size=11, color=SLATE, align=PP_ALIGN.CENTER)
    text(s, Inches(10.95), Inches(4.0), Inches(1.75), Inches(0.4), "< 1 sec", size=11, color=SLATE, align=PP_ALIGN.CENTER)
    rect(s, Inches(11.05), Inches(4.55), Inches(1.55), Inches(0.02), LINE)
    text(s, Inches(10.95), Inches(4.65), Inches(1.75), Inches(0.4), "vs 2 ore", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, Inches(10.95), Inches(4.9), Inches(1.75), Inches(0.4), "Excel manuale", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    for i in range(4):
        y = Inches(2.7) + i * Inches(0.95)
        arrow(s, Inches(3.4), y, Inches(4.3), Inches(4.25), color=MUTED)
    for i in range(3):
        y = Inches(2.92) + i * Inches(1.0)
        arrow(s, Inches(6.8), Inches(4.25), Inches(7.8), y, color=MUTED)

    rect(s, Inches(0.6), Inches(5.85), Inches(12.13), Inches(0.9), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.85), Inches(5.95), Inches(3.5), Inches(0.3), "5 LIVELLI DI PRIORITÀ", size=10, bold=True, color=GOLD)
    crits = ["Urgenza", "Ritardo", "Affinità batch", "Sequenza ciclo", "Formato"]
    for i, c in enumerate(crits):
        x = Inches(0.85) + Inches(2.4) * i
        text(s, x, Inches(6.3), Inches(2.3), Inches(0.4), f"0{i+1}  {c}", size=12, bold=True, color=WHITE)

    footer(s, 8, total)

    # 9 Priorità — 5 criteri spiegati
    s = blank(p)
    header(s, "Funzione di scoring", "Cinque criteri ponderati, sommati in un solo punteggio per fase", num="02")
    levels = [
        ("01", "URGENZA CONTRATTUALE", "+1000",
         "Distanza in giorni dalla data di consegna concordata.",
         "Criticità ≤ 3 giorni · peso dominante sull'intera funzione", RED),
        ("02", "DEVIAZIONE DI PIANO", "+500",
         "Scostamento tra inizio fase atteso e orologio reale.",
         "Compensazione automatica del ritardo accumulato a monte", AMBER),
        ("03", "AFFINITÀ DI BATCH", "+200",
         "Compatibilità di setup con commesse adiacenti in coda.",
         "Stesso colore, fustella, lastra · setup time prossimo a zero", BLUE),
        ("04", "DIPENDENZE DI CICLO", "+100",
         "Pronta-disponibilità della fase a monte come prerequisito.",
         "Vincoli di precedenza rispettati · zero attese cieche", TEAL),
        ("05", "AFFINITÀ DI FORMATO", "+50",
         "Continuità di formato foglio sulla stessa unità produttiva.",
         "Riduzione delle tarature · throughput più costante", GREEN),
    ]
    for i, (num, name, pts, desc, sub, c) in enumerate(levels):
        y = Inches(2.0) + i * Inches(1.0)
        rect(s, Inches(0.6), y, Inches(12.13), Inches(0.9), WHITE, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, Inches(0.6), y, Inches(0.08), Inches(0.9), c)
        text(s, Inches(0.85), y + Inches(0.15), Inches(0.7), Inches(0.6), num, size=22, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.6), y + Inches(0.08), Inches(3.1), Inches(0.4), name, size=12, bold=True, color=INK)
        text(s, Inches(1.6), y + Inches(0.43), Inches(3.1), Inches(0.45), pts + " punti", size=11, bold=True, color=c)
        text(s, Inches(4.85), y + Inches(0.1), Inches(7.7), Inches(0.4), desc, size=12, color=INK)
        text(s, Inches(4.85), y + Inches(0.48), Inches(7.7), Inches(0.4), sub, size=10, color=SLATE)
    footer(s, 9, total)

    # 10 Esempio scoring
    s = blank(p)
    header(s, "Caso applicato", "Una commessa reale · cinque criteri · un punteggio · una posizione in coda", num="02")

    rect(s, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.9), WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), Inches(2.0), Inches(4.5), Inches(0.5), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.6), Inches(2.0), Inches(4.5), Inches(0.5), "COMMESSA ESEMPIO", size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.85), Inches(2.7), Inches(4), Inches(0.5), "Astuccio packaging", size=18, bold=True, color=INK)
    text(s, Inches(0.85), Inches(3.05), Inches(4), Inches(0.4), "Settore alimentare premium", size=12, color=SLATE)
    rect(s, Inches(0.85), Inches(3.55), Inches(4), Inches(0.02), LINE)
    info = [
        ("Quantità", "13.800 fg"),
        ("Lead time residuo", "2 giorni"),
        ("Unità produttiva", "linea offset principale"),
        ("Configurazione", "4C + verniciatura"),
        ("Stato a monte", "fase precedente pronta"),
        ("Job adiacente", "stessa configurazione"),
    ]
    for i, (k, v) in enumerate(info):
        y = Inches(3.7) + i * Inches(0.5)
        text(s, Inches(0.85), y, Inches(2), Inches(0.4), k, size=11, color=MUTED)
        text(s, Inches(2.85), y, Inches(2), Inches(0.4), v, size=12, bold=True, color=INK)

    text(s, Inches(5.4), Inches(2.0), Inches(4), Inches(0.4), "PUNTEGGIO PER CRITERIO", size=11, bold=True, color=BLUE)
    scoring = [
        ("Urgenza · ≤ 3 giorni", "+1000", BLUE, 0.95),
        ("Ritardo · in tempo", "0", LINE, 0.0),
        ("Affinità batch · stesso 4C", "+200", BLUE, 0.20),
        ("Sequenza · fase pronta", "+100", BLUE, 0.10),
        ("Formato · invariato", "+50", BLUE, 0.05),
    ]
    bar_x = Inches(8.3)
    bar_w_max = Inches(3.8)
    for i, (label, pts, c, frac) in enumerate(scoring):
        y = Inches(2.5) + i * Inches(0.7)
        text(s, Inches(5.4), y + Inches(0.1), Inches(2.7), Inches(0.5), label, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, bar_x, y + Inches(0.18), bar_w_max, Inches(0.28), BG)
        if frac > 0:
            bw = Emu(int(bar_w_max * frac))
            rect(s, bar_x, y + Inches(0.18), bw, Inches(0.28), c)
        text(s, bar_x + bar_w_max + Inches(0.15), y + Inches(0.1), Inches(0.7), Inches(0.5), pts, size=12, bold=True, color=c if frac > 0 else MUTED, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    rect(s, Inches(5.4), Inches(6.05), Inches(7.3), Inches(0.02), INK)
    text(s, Inches(5.4), Inches(6.2), Inches(3), Inches(0.5), "TOTALE", size=14, bold=True, color=INK)
    text(s, Inches(10.7), Inches(6.15), Inches(2), Inches(0.55), "1.350", size=28, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    rect(s, Inches(0.6), Inches(7.0), Inches(12.13), Inches(0.04), GOLD)
    text(s, Inches(0.6), Inches(7.1), Inches(12.13), Inches(0.3), "Score 1.350 → testa coda · ricalcolo automatico a ogni evento di pianificazione", size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 10, total)

    # 11 Sect 03
    section_divider(p, "03", "Impatto operativo", "Throughput, lead time, OEE.\nMisurati. Non promessi.")

    # 12 KPI
    s = blank(p)
    header(s, "Otto KPI di riferimento", "Baseline misurata su sito attivo · benchmark di settore disponibile in audit", num="03")
    rect(s, 0, Inches(2.0), Inches(13.333), Inches(5), BG)
    kpi(s, Inches(0.6), Inches(2.3), Inches(2.95), Inches(2.1), "< 1 s", "Latenza ricomposizione 50+ commesse", color=BLUE)
    kpi(s, Inches(3.7), Inches(2.3), Inches(2.95), Inches(2.1), "−95%", "Effort scheduling capo reparto", color=GREEN)
    kpi(s, Inches(6.8), Inches(2.3), Inches(2.95), Inches(2.1), "−20%", "Setup time complessivo", color=GREEN)
    kpi(s, Inches(9.9), Inches(2.3), Inches(2.85), Inches(2.1), "+15%", "OEE attesa, stampa offset", color=GREEN)
    kpi(s, Inches(0.6), Inches(4.6), Inches(2.95), Inches(2.1), "5", "Criteri di scoring ponderati", color=NAVY)
    kpi(s, Inches(3.7), Inches(4.6), Inches(2.95), Inches(2.1), "Auto", "Batching · alert colli di bottiglia", color=AMBER)
    kpi(s, Inches(6.8), Inches(4.6), Inches(2.95), Inches(2.1), "Real-time", "Gantt distribuito, sempre allineato", color=BLUE)
    kpi(s, Inches(9.9), Inches(4.6), Inches(2.85), Inches(2.1), "Mobile-first", "Tablet a bordo macchina", color=BLUE)
    footer(s, 12, total)

    # 13 Benefici operativi
    s = blank(p)
    header(s, "Quattro cambiamenti misurabili", "Confronto baseline vs steady-state · entro le prime quattro settimane", num="03")
    benefits = [
        ("Effort di scheduling", "Pianificazione delegata al motore. Capo reparto liberato.",
         "BASELINE  ~2 h / die         POST  ~15 min / die", BLUE),
        ("Reattività al cambio", "Riprogrammazione automatica a ogni evento di reparto.",
         "BASELINE  30-60 min          POST  ~ 1 secondo", GREEN),
        ("Setup time", "Sequenze a setup compatibile, batching automatico.",
         "BASELINE  setup ripetuti     POST  −20% tempo unità produttiva", AMBER),
        ("Visibilità operativa", "Stesso Gantt per capo reparto, operatori, direzione.",
         "BASELINE  voce + Excel       POST  Gantt vivo, distribuito", TEAL),
    ]
    for i, (n, d, comp, c) in enumerate(benefits):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(6.18)
        y = Inches(2.0) + row * Inches(2.45)
        rect(s, x, y, Inches(5.95), Inches(2.25), WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, Inches(0.08), Inches(2.25), c)
        text(s, x + Inches(0.3), y + Inches(0.2), Inches(5.5), Inches(0.5), n, size=18, bold=True, color=INK)
        text(s, x + Inches(0.3), y + Inches(0.75), Inches(5.5), Inches(0.8), d, size=12, color=SLATE, spacing=1.3)
        rect(s, x + Inches(0.3), y + Inches(1.5), Inches(5.5), Inches(0.02), LINE)
        text(s, x + Inches(0.3), y + Inches(1.6), Inches(5.5), Inches(0.5), comp, size=11, bold=True, color=c)
    footer(s, 13, total)

    # 14 Sect 04
    section_divider(p, "04", "Adozione", "Quattro settimane al go-live.\nUna garanzia formale.")

    # 15 Timeline
    s = blank(p)
    header(s, "Roadmap di adozione · 4 settimane", "Dalla firma del contratto alla prima commessa pianificata da Mossa 37", num="04")
    weeks = [
        ("S1", "Discovery", "Audit del flusso · mapping unità produttive · cicli di lavorazione", BLUE),
        ("S2", "Setup", "Installazione · connettori · configurazione regole di scoring", TEAL),
        ("S3", "Pilot", "Training capo reparto + operatori · settimana di scheduling assistito", AMBER),
        ("S4", "Go-live", "Affiancamento on-site · tuning parametri · handover formale", GREEN),
    ]
    rect(s, Inches(0.6), Inches(3.5), Inches(12.13), Inches(0.06), LINE)
    cw = Inches(2.95)
    gx = Inches(0.13)
    for i, (w, n, d, c) in enumerate(weeks):
        x = Inches(0.6) + i * (cw + gx)
        rect(s, x + cw / 2 - Inches(0.12), Inches(3.4), Inches(0.24), Inches(0.24), c, shape=MSO_SHAPE.OVAL)
        rect(s, x, Inches(2.0), cw, Inches(1.3), WHITE, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(2.0), cw, Inches(0.4), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x, Inches(2.0), cw, Inches(0.4), w + "  ·  Settimana " + str(i + 1), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(0.2), Inches(2.5), cw - Inches(0.3), Inches(0.5), n, size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.2), Inches(2.95), cw - Inches(0.3), Inches(0.4), d, size=10, color=SLATE, align=PP_ALIGN.CENTER)
        rect(s, x, Inches(3.85), cw, Inches(1.6), BG, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x + Inches(0.2), Inches(3.95), cw - Inches(0.3), Inches(0.4), "DELIVERABLE", size=10, bold=True, color=c)
        deliverables = [
            "Documento as-is + cicli\nMappa macchine\nLista priorità",
            "Ambiente attivo\nConnettori configurati\nUtenti creati",
            "Dashboard live\nManuale operatore\n2 sessioni training",
            "Go-live ufficiale\nReport prima settimana\n90gg supporto",
        ][i]
        text(s, x + Inches(0.2), Inches(4.25), cw - Inches(0.3), Inches(1.2), deliverables, size=11, color=INK, spacing=1.3)

    rect(s, Inches(0.6), Inches(5.85), Inches(12.13), Inches(1.1), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.4), "GARANZIA TEMPI", size=11, bold=True, color=GOLD)
    text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6), "Se non andiamo in produzione entro la 4ª settimana, il primo mese è gratuito.", size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 15, total)

    # 16 Garanzie
    s = blank(p)
    header(s, "Quattro garanzie contrattuali", "Impegni formalizzati in clausole, non slide marketing", num="04")
    guarantees = [
        ("30 giorni", "Diritto di recesso pieno nei primi 30 giorni post go-live.", "Right to exit"),
        ("99,5%", "SLA uptime piattaforma. Sotto soglia · credito automatico.", "Affidabilità"),
        ("90 giorni", "Supporto premium incluso oltre il go-live, senza maggiorazioni.", "Hyper-care"),
        ("Sempre", "Aggiornamenti continui inclusi nel canone, senza upgrade fee.", "Continuous delivery"),
    ]
    cw = Inches(2.95)
    gx = Inches(0.13)
    for i, (val, d, tag) in enumerate(guarantees):
        x = Inches(0.6) + i * (cw + gx)
        rect(s, x, Inches(2.0), cw, Inches(3.5), WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(2.0), cw, Inches(0.06), GOLD)
        text(s, x + Inches(0.3), Inches(2.25), cw - Inches(0.4), Inches(0.4), tag.upper(), size=11, bold=True, color=BLUE)
        text(s, x + Inches(0.3), Inches(2.7), cw - Inches(0.4), Inches(1.0), val, size=36, bold=True, color=INK)
        text(s, x + Inches(0.3), Inches(3.7), cw - Inches(0.4), Inches(1.6), d, size=12, color=SLATE, spacing=1.3)

    rect(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(1.3), BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.4), "CANALI DI SUPPORTO", size=11, bold=True, color=BLUE)
    chans = [("Email", "risposta 4h"), ("Telefono", "8-18 lun-ven"), ("Chat", "in app"), ("On-site", "su richiesta")]
    for i, (n, d) in enumerate(chans):
        x = Inches(0.9) + i * Inches(2.95)
        text(s, x, Inches(6.25), Inches(2.8), Inches(0.4), n, size=15, bold=True, color=INK)
        text(s, x, Inches(6.6), Inches(2.8), Inches(0.4), d, size=11, color=SLATE)
    footer(s, 16, total)

    # 17 — Trigger di ricomposizione (event-driven)
    s = blank(p)
    header(s, "Trigger di ricomposizione", "Sei eventi · sei riprogrammazioni automatiche · zero intervento manuale", num="02")
    triggers = [
        ("01", "Fase terminata", "L'operatore chiude → il motore propaga la fase successiva", BLUE),
        ("02", "Fase avviata in ritardo", "Il piano si riassesta sui ritardi reali, non su quelli stimati", AMBER),
        ("03", "Variazione urgenza", "Il cliente sposta una consegna → riprogrammazione in 1 secondo", RED),
        ("04", "Nuova commessa da ERP", "L'ordine entra → posizionato in coda secondo i criteri", GREEN),
        ("05", "Indisponibilità macchina", "Fermo o manutenzione → carico ridistribuito automaticamente", TEAL),
        ("06", "Cambio finestra produttiva", "Apertura/chiusura reparto, nuova squadra → risorse ricalcolate", NAVY),
    ]
    cw = Inches(4.05)
    ch = Inches(2.3)
    gx = Inches(0.13)
    gy = Inches(0.2)
    x0 = Inches(0.6)
    y0 = Inches(2.0)
    for i, (num, n, d, c) in enumerate(triggers):
        col = i % 3
        row = i // 3
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, Inches(0.08), ch, c)
        text(s, x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.4), num, size=11, bold=True, color=c)
        text(s, x + Inches(0.3), y + Inches(0.55), cw - Inches(0.4), Inches(0.6), n, size=18, bold=True, color=INK)
        text(s, x + Inches(0.3), y + Inches(1.2), cw - Inches(0.4), Inches(1), d, size=12, color=SLATE, spacing=1.3)
    footer(s, 17, total)

    # 18 — Configurabilità
    s = blank(p)
    header(s, "Cosa potete configurare", "Le regole vivono dove vivono i processi · zero codice", num="02")
    configs = [
        ("Calendari & disponibilità", "Orari macchina, festività, turni se presenti, pause programmate. Anche senza sistema presenze esterno."),
        ("Pesi dei criteri", "Tarate i cinque livelli di priorità sul vostro modello di business."),
        ("Soglie di alert", "Sopra-soglia di carico, lead time critico, ritardo accumulato sulla fase."),
        ("Affinità di setup", "Definite cosa rende due commesse 'consecutive' sulla vostra unità produttiva."),
        ("Vincoli di precedenza", "Sequenza obbligatoria delle fasi, eccezioni per cicli speciali."),
        ("Forzature manuali", "Il capo reparto può sovrascrivere lo score quando serve. Sempre tracciato."),
    ]
    for i, (n, d) in enumerate(configs):
        y = Inches(2.0) + i * Inches(0.78)
        rect(s, Inches(0.6), y, Inches(12.13), Inches(0.7), WHITE, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, Inches(0.6), y, Inches(0.08), Inches(0.7), GOLD)
        text(s, Inches(0.95), y + Inches(0.13), Inches(3.8), Inches(0.5), n, size=14, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(4.85), y + Inches(0.15), Inches(7.8), Inches(0.45), d, size=12, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.6), Inches(6.7), Inches(12.13), Inches(0.04), GOLD)
    text(s, Inches(0.6), Inches(6.78), Inches(12), Inches(0.3), "Tutta la configurazione gestita in interfaccia owner · nessun developer richiesto", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, 18, total)

    closing_slide(p,
        "Mezza giornata. Zero impegno.\nUn piano scritto.",
        "Audit on-site del vostro flusso reale.\nPoi decidete con i numeri in mano, non con la fiducia.",
        "Contatti su richiesta")

    out = os.path.join(DOCS, "Mossa37_Scheduler_Pitch.pptx")
    p.save(out)
    return out


# ============================================================
# PPT 2 — MES Overview
# ============================================================
def build_mes_overview():
    p = new_pres()
    total = 16

    hero_title(p,
        "MES  ·  MANUFACTURING EXECUTION SYSTEM",
        "Una sola piattaforma.\nTutta la produzione.",
        "Dal preventivo alla bolla in uscita: orchestrazione end-to-end,\ndato unico, decisioni in tempo reale.",
        "Documento riservato  ·  Aprile 2026")

    # 2 Indice
    s = blank(p)
    header(s, "Agenda", "Quindici minuti. Cinque sezioni. Una decisione strategica.", num="00")
    items = [
        ("01", "Frammentazione operativa", "perché ERP, macchine, Excel non sono un sistema"),
        ("02", "Il MES come orchestratore", "moduli, integrazioni, modello dati unificato"),
        ("03", "Impatto operativo", "throughput, lead time, ore-uomo, OEE"),
        ("04", "Roadmap di adozione", "quattro settimane al go-live, garanzia formale"),
        ("05", "Demo live", "la piattaforma in funzione, sui vostri scenari"),
    ]
    for i, (n, t, d) in enumerate(items):
        y = Inches(2.0) + i * Inches(0.95)
        rect(s, Inches(0.6), y, Inches(12.13), Inches(0.85), WHITE, line=LINE, line_w=Pt(0.5))
        text(s, Inches(0.9), y + Inches(0.18), Inches(0.8), Inches(0.5), n, size=22, bold=True, color=BLUE)
        text(s, Inches(2.0), y + Inches(0.18), Inches(4.5), Inches(0.5), t, size=18, bold=True, color=INK)
        text(s, Inches(6.7), y + Inches(0.22), Inches(6), Inches(0.5), d, size=13, color=SLATE)
    footer(s, 2, total)

    # 3 Sect 01
    section_divider(p, "01", "Frammentazione", "Sistemi isolati.\nDecisioni a memoria.")

    # 4 Sintomi
    s = blank(p)
    header(s, "Cinque sintomi diagnostici", "Se ne riconoscete almeno tre, il vostro stack non è un sistema", num="01")
    syms = [
        ("ERP silos", "Il gestionale registra le transazioni. Il dato non viaggia oltre."),
        ("Excel sprawl", "Pianificazione, ore, scarti, spedizioni: dieci file, dieci verità."),
        ("Macchine offline", "Le linee produttive generano telemetria. Nessuno la consuma."),
        ("Operatori senza interfaccia", "Foglio stampato e penna. Decisioni operative reattive."),
        ("Direzione senza dashboard", "Reportistica a fine mese. Decisioni necessarie real-time."),
    ]
    for i, (n, d) in enumerate(syms):
        y = Inches(2.0) + i * Inches(0.95)
        rect(s, Inches(0.6), y, Inches(12.13), Inches(0.85), WHITE, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, Inches(0.6), y, Inches(0.08), Inches(0.85), RED)
        text(s, Inches(0.95), y + Inches(0.2), Inches(3.5), Inches(0.5), n, size=16, bold=True, color=INK)
        text(s, Inches(4.5), y + Inches(0.22), Inches(8), Inches(0.5), d, size=13, color=SLATE)
    footer(s, 4, total)

    # 5 Sect 02
    section_divider(p, "02", "Orchestrazione", "Un solo modello dati.\nUna sola interfaccia.")

    # 6 Manifesto + pilastri
    s = blank(p)
    header(s, "Il MES, in profondità", "Manufacturing Execution System verticale per packaging cartotecnico", num="02")
    rect(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.2), BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2),
         "Orchestrazione end-to-end del flusso produttivo: preventivo, scheduling,\nesecuzione di reparto, magazzino, spedizione, KPI direzionali.\nUn solo dato vero. In tempo reale. Su tutta la filiera.",
         size=20, bold=True, color=INK, spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)
    pillars = [
        ("Connesso", "integrazioni native con il vostro stack esistente"),
        ("Real-time", "modello dati event-driven, latenza sub-secondo"),
        ("Mobile-first", "operatori, magazzino, capo reparto, da qualunque device"),
        ("Verticale", "logiche calibrate sulla cartotecnica, non generiche"),
    ]
    cw = Inches(2.95)
    for i, (n, d) in enumerate(pillars):
        x = Inches(0.6) + i * (cw + Inches(0.13))
        rect(s, x, Inches(4.5), cw, Inches(2.4), WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(4.5), Inches(0.08), Inches(2.4), GOLD)
        text(s, x + Inches(0.3), Inches(4.8), cw - Inches(0.4), Inches(0.6), n, size=22, bold=True, color=INK)
        text(s, x + Inches(0.3), Inches(5.5), cw - Inches(0.4), Inches(1.2), d, size=12, color=SLATE, spacing=1.3)
    footer(s, 6, total)

    # 7 Moduli
    s = blank(p)
    header(s, "Architettura modulare · sei capability", "Moduli indipendenti · modello dati condiviso · adozione progressiva", num="02")
    cols = [
        ("01", "Shop floor", "Esecuzione di reparto · avanzamento fase · annotazioni operative", BLUE),
        ("02", "Direzione", "Cruscotto KPI · ordini · ore-uomo · costo industriale", NAVY),
        ("03", "Prestampa", "Lavorazioni · matrici · fustelle · note tecniche", AMBER),
        ("04", "Spedizione", "Bolle in uscita · integrazione corrieri · tracking · esiti", GREEN),
        ("05", "Magazzino", "Tracciabilità materia prima · QR · ricezione AI-assisted", TEAL),
        ("06", "Scheduler", "Mossa 37 · scheduling real-time · cinque criteri di priorità", RED),
    ]
    cw = Inches(4.05)
    ch = Inches(2.3)
    gx = Inches(0.13)
    gy = Inches(0.2)
    x0 = Inches(0.6)
    y0 = Inches(2.0)
    for i, (num, n, d, c) in enumerate(cols):
        col = i % 3
        row = i // 3
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, Inches(0.08), ch, c)
        text(s, x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.4), num, size=11, bold=True, color=c)
        text(s, x + Inches(0.3), y + Inches(0.55), cw - Inches(0.4), Inches(0.6), n, size=22, bold=True, color=INK)
        text(s, x + Inches(0.3), y + Inches(1.2), cw - Inches(0.4), Inches(1), d, size=12, color=SLATE, spacing=1.3)
    footer(s, 7, total)

    # 8 Integrazioni
    s = blank(p)
    header(s, "Otto pattern di integrazione", "Architettura non invasiva · si appoggia allo stack esistente", num="02")
    integ = [
        ("ERP gestionale", "ordini · commesse · clienti · priorità", BLUE),
        ("Macchine offset", "telemetria di produzione · scarti · avanzamento", NAVY),
        ("Macchine digitali", "coda di stampa · job · fogli prodotti", AMBER),
        ("Sistema presenze", "opzionale · integriamo il vostro o gestione interna", GREEN),
        ("Corrieri", "spedizioni · tracking · esiti consegna", RED),
        ("Mobile / chat", "operatori in mobilità · notifiche push", BLUE),
        ("AI vision", "OCR documenti · classificazione · ispezione", TEAL),
        ("API custom", "qualsiasi sorgente proprietaria · connettore su misura", GREEN),
    ]
    cw = Inches(2.95)
    ch = Inches(1.4)
    for i, (n, d, c) in enumerate(integ):
        col = i % 4
        row = i // 4
        x = Inches(0.6) + col * (cw + Inches(0.13))
        y = Inches(2.0) + row * (ch + Inches(0.2))
        rect(s, x, y, cw, ch, WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, cw, Inches(0.4), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x, y, cw, Inches(0.4), n, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(0.2), y + Inches(0.55), cw - Inches(0.3), Inches(0.85), d, size=11, color=SLATE)
    footer(s, 8, total)

    # 9 Sect 03
    section_divider(p, "03", "Impatto operativo", "Throughput, lead time, OEE.\nMisurati. Non promessi.")

    # 10 KPI
    s = blank(p)
    header(s, "Otto KPI di riferimento", "Baseline misurata su sito attivo · benchmark di settore disponibile in audit", num="03")
    rect(s, 0, Inches(2.0), Inches(13.333), Inches(5), BG)
    kpi(s, Inches(0.6), Inches(2.3), Inches(2.95), Inches(2.1), "0", "Spreadsheet di passaggio", color=GREEN)
    kpi(s, Inches(3.7), Inches(2.3), Inches(2.95), Inches(2.1), "< 1 s", "Latenza ricomposizione scheduler", color=BLUE)
    kpi(s, Inches(6.8), Inches(2.3), Inches(2.95), Inches(2.1), "5", "Stati di avanzamento fase", color=NAVY)
    kpi(s, Inches(9.9), Inches(2.3), Inches(2.85), Inches(2.1), "Real-time", "Cruscotto direzionale, sempre vivo", color=BLUE)
    kpi(s, Inches(0.6), Inches(4.6), Inches(2.95), Inches(2.1), "Auto", "Ingresso merce AI-assisted", color=AMBER)
    kpi(s, Inches(3.7), Inches(4.6), Inches(2.95), Inches(2.1), "QR", "Tracciabilità unitaria sui bancali", color=TEAL)
    kpi(s, Inches(6.8), Inches(4.6), Inches(2.95), Inches(2.1), "Mobile", "Tablet di reparto + mobile magazzino", color=BLUE)
    kpi(s, Inches(9.9), Inches(4.6), Inches(2.85), Inches(2.1), "1 s", "Polling notifiche operative", color=NAVY)
    footer(s, 10, total)

    # 11 Magazzino AI workflow
    s = blank(p)
    header(s, "Caso applicato: magazzino AI-assisted", "Ricezione merce da foto bolla a giacenza in meno di un minuto", num="03")
    wf = [
        ("Foto bolla", "fornitore"),
        ("Chat / mobile", "bot riceve"),
        ("AI vision", "estrae dati"),
        ("Auto-carico", "se in anagrafica"),
        ("Etichetta QR", "stampa bancale"),
    ]
    nw = Inches(2.3)
    nh = Inches(1.1)
    gap = Inches(0.15)
    x0 = Inches(0.6)
    y0 = Inches(2.1)
    for i, (n, d) in enumerate(wf):
        x = x0 + i * (nw + gap)
        node(s, x, y0, nw, nh, n, sub=d, accent=TEAL)
        if i < len(wf) - 1:
            arrow(s, x + nw, y0 + Inches(0.55), x + nw + gap, y0 + Inches(0.55), color=TEAL, weight=1.5)

    bullets(s, Inches(0.6), Inches(3.7), Inches(7.5), Inches(2.8), [
        "Acquisizione: foto bolla via mobile, OCR vision-AI",
        "Estrazione strutturata: codice articolo, fogli, kg, bancali, lotto",
        "Carico automatico: creazione movimento + aggiornamento giacenza",
        "Etichettatura: QR generato per ogni bancale, tracciabilità unitaria",
        "Consumo: scarico in fase stampa, validato server-side",
    ], size=14)
    # right side panel — micro KPIs
    rect(s, Inches(8.3), Inches(3.7), Inches(4.43), Inches(2.85), BG, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(8.5), Inches(3.85), Inches(4), Inches(0.4), "DELTA OPERATIVO", size=11, bold=True, color=TEAL)
    micros = [
        ("Data entry", "manuale → automatico"),
        ("Tempo ciclo ricezione", "5 min → 30 sec"),
        ("Errori di trascrizione", "azzerati"),
    ]
    for i, (k, v) in enumerate(micros):
        y = Inches(4.3) + i * Inches(0.75)
        text(s, Inches(8.5), y, Inches(4), Inches(0.4), k, size=12, bold=True, color=INK)
        text(s, Inches(8.5), y + Inches(0.35), Inches(4), Inches(0.4), v, size=11, color=SLATE)
    # bottom strip
    rect(s, Inches(0.6), Inches(6.7), Inches(12.13), Inches(0.04), GOLD)
    text(s, Inches(0.6), Inches(6.78), Inches(12), Inches(0.3), "Audit trail completo · chi · cosa · quando · da quale ordine · pronto per certificazioni", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, 11, total)

    # 12 Sect 04
    section_divider(p, "04", "Adozione", "Quattro settimane al go-live.\nUna garanzia formale.")

    # 13 Timeline
    s = blank(p)
    header(s, "Roadmap di adozione · 4 settimane", "Dalla firma del contratto al primo bancale tracciato", num="04")
    weeks = [
        ("S1", "Discovery", "Audit del flusso · mapping capability · cicli produttivi", BLUE),
        ("S2", "Setup", "Provisioning ambiente · connettori · utenti · permessi", TEAL),
        ("S3", "Pilot", "Training shop floor + direzione · settimana di esecuzione assistita", AMBER),
        ("S4", "Go-live", "Affiancamento on-site · tuning · handover formale", GREEN),
    ]
    rect(s, Inches(0.6), Inches(3.5), Inches(12.13), Inches(0.06), LINE)
    cw = Inches(2.95)
    gx = Inches(0.13)
    for i, (w, n, d, c) in enumerate(weeks):
        x = Inches(0.6) + i * (cw + gx)
        rect(s, x + cw / 2 - Inches(0.12), Inches(3.4), Inches(0.24), Inches(0.24), c, shape=MSO_SHAPE.OVAL)
        rect(s, x, Inches(2.0), cw, Inches(1.3), WHITE, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(2.0), cw, Inches(0.4), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x, Inches(2.0), cw, Inches(0.4), w + "  ·  Settimana " + str(i + 1), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(0.2), Inches(2.5), cw - Inches(0.3), Inches(0.5), n, size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.2), Inches(2.95), cw - Inches(0.3), Inches(0.4), d, size=10, color=SLATE, align=PP_ALIGN.CENTER)
        rect(s, x, Inches(3.85), cw, Inches(1.6), BG, line=LINE, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x + Inches(0.2), Inches(3.95), cw - Inches(0.3), Inches(0.4), "DELIVERABLE", size=10, bold=True, color=c)
        deliverables = [
            "Documento as-is\nMappa moduli\nCicli produttivi",
            "Ambiente live\nConnettori attivi\nUtenti creati",
            "Dashboard testate\nManuale operatore\n3 sessioni training",
            "Go-live ufficiale\nReport prima settimana\n90gg supporto",
        ][i]
        text(s, x + Inches(0.2), Inches(4.25), cw - Inches(0.3), Inches(1.2), deliverables, size=11, color=INK, spacing=1.3)
    rect(s, Inches(0.6), Inches(5.85), Inches(12.13), Inches(1.1), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.4), "GARANZIA TEMPI", size=11, bold=True, color=GOLD)
    text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6), "Se non andiamo in produzione entro la 4ª settimana, il primo mese è gratuito.", size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 13, total)

    # 14 — Modello dati unificato
    s = blank(p)
    header(s, "Un solo modello dati", "Sei capability che condividono le stesse entità · zero sincronizzazioni", num="04")
    entities = [
        ("Cliente", "anagrafica · ordini · cronologia consegne", BLUE),
        ("Commessa", "ciclo · priorità · stato · scadenze", NAVY),
        ("Fase", "macchina · operatore · tempi · scarti", AMBER),
        ("Operatore", "skill · disponibilità · permessi · audit", GREEN),
        ("Articolo magazzino", "giacenza · lotti · QR · movimenti", TEAL),
        ("Spedizione", "bolla · corriere · tracking · esito", RED),
    ]
    cw = Inches(4.05)
    ch = Inches(2.3)
    gx = Inches(0.13)
    gy = Inches(0.2)
    x0 = Inches(0.6)
    y0 = Inches(2.0)
    for i, (n, d, c) in enumerate(entities):
        col = i % 3
        row = i // 3
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, Inches(0.08), ch, c)
        text(s, x + Inches(0.3), y + Inches(0.4), cw - Inches(0.4), Inches(0.6), n, size=22, bold=True, color=INK)
        text(s, x + Inches(0.3), y + Inches(1.1), cw - Inches(0.4), Inches(1.1), d, size=12, color=SLATE, spacing=1.3)
    footer(s, 14, total)

    # 15 — Sicurezza & compliance
    s = blank(p)
    header(s, "Sicurezza, audit, compliance", "Quattro pilastri · pronti per audit cliente, ISO, GDPR", num="04")
    pillars = [
        ("Audit trail", "Ogni evento tracciato: chi, cosa, quando, da quale ordine. Immutabile.", "Tracciabilità"),
        ("Permessi RBAC", "Ruoli granulari: operatore vede solo il suo reparto, owner vede tutto.", "Accesso controllato"),
        ("GDPR-ready", "Registro trattamenti, retention configurabile, diritto all'oblio nativo.", "Dati personali"),
        ("Backup & DR", "Snapshot quotidiani, RPO < 24h, RTO < 4h. Recovery testato.", "Continuità"),
    ]
    cw = Inches(2.95)
    gx = Inches(0.13)
    for i, (n, d, tag) in enumerate(pillars):
        x = Inches(0.6) + i * (cw + gx)
        rect(s, x, Inches(2.0), cw, Inches(4.6), WHITE, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(2.0), cw, Inches(0.06), GOLD)
        text(s, x + Inches(0.3), Inches(2.25), cw - Inches(0.4), Inches(0.4), tag.upper(), size=10, bold=True, color=BLUE)
        text(s, x + Inches(0.3), Inches(2.7), cw - Inches(0.4), Inches(0.7), n, size=20, bold=True, color=INK)
        text(s, x + Inches(0.3), Inches(3.6), cw - Inches(0.4), Inches(2.8), d, size=12, color=SLATE, spacing=1.4)
    rect(s, Inches(0.6), Inches(7.0), Inches(12.13), Inches(0.04), GOLD)
    text(s, Inches(0.6), Inches(7.08), Inches(12), Inches(0.3), "Documentazione completa fornita: registro trattamenti · DPIA · piano sicurezza · runbook", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, 15, total)

    # 16 Closing
    closing_slide(p,
        "Mezza giornata. Zero impegno.\nUn piano scritto.",
        "Audit on-site del vostro flusso reale.\nPoi decidete con i numeri in mano, non con la fiducia.",
        "Contatti su richiesta")

    out = os.path.join(DOCS, "MES_GraficaNappa_Overview.pptx")
    p.save(out)
    return out


if __name__ == "__main__":
    a = build_scheduler()
    b = build_mes_overview()
    print("OK:", a)
    print("OK:", b)
